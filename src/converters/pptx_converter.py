"""Conversor PPTX -> texto. Story 1.6 — python-pptx + fallback OCR pytesseract."""

from __future__ import annotations

import io
import re
import unicodedata
from concurrent.futures import ThreadPoolExecutor, as_completed
from contextlib import contextmanager
from pathlib import Path
from typing import Any

import pytesseract
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


@contextmanager
def _open_Presentation(path: str):  # noqa: N802
    """Context manager para Presentation (M1 — python-pptx não tem __exit__ nativo)."""
    prs = Presentation(path)
    try:
        yield prs
    finally:
        del prs


MAX_FILE_BYTES = 10 * 1024 * 1024  # 10 MB
MIN_CHARS_SLIDE = 10  # M5 — threshold por slide para acionar OCR
MAX_SLIDES = 200  # M5 — limite de slides para evitar processamento excessivo
MAX_WORKERS_BATCH = 4  # R5 — paralelismo no batch de 120 arquivos


def convert_pptx(file_path: str | Path) -> dict[str, object]:
    """Converte um arquivo PPTX em texto.

    Returns:
        dict com campos:
        - texto (str): texto extraído concatenado de todos os slides
        - slides (int): número de slides
        - metodo_extracao (str): "nativo", "ocr" ou "misto"
        - flags (list[str]): vazio = sucesso; possíveis valores:
            "arquivo_muito_grande", "arquivo_ilegivel", "formato_nao_suportado"
    """
    file_path = Path(file_path)

    # R4 — rejeitar .ppt legado antes de tentar Presentation()
    if file_path.suffix.lower() == ".ppt":
        return _falha("formato_nao_suportado", slides=0)

    try:
        tamanho = file_path.stat().st_size
    except OSError:
        return _falha("arquivo_ilegivel", slides=0)

    if tamanho > MAX_FILE_BYTES:
        mb = tamanho // (1024 * 1024)
        raise ValueError(f"Arquivo muito grande: {mb}MB. Limite: 10MB.")

    try:
        with _open_Presentation(str(file_path)) as prs:  # M1 — context manager
            total_slides = len(prs.slides)
            textos: list[str] = []
            metodos: set[str] = set()

            for slide in list(prs.slides)[:MAX_SLIDES]:  # M5 — limite de slides
                texto_slide = _extrair_text_frames(slide)

                if len(texto_slide.strip()) < MIN_CHARS_SLIDE and _tem_imagem(slide.shapes):
                    texto_slide = _ocr_slide(slide)
                    metodos.add("ocr")
                else:
                    metodos.add("nativo")

                textos.append(texto_slide)

            if total_slides > MAX_SLIDES:
                total_slides = MAX_SLIDES

            texto_final = "\n\n".join(t for t in textos if t.strip())

            if "ocr" in metodos and "nativo" in metodos:
                metodo = "misto"
            elif "ocr" in metodos:
                metodo = "ocr"
            else:
                metodo = "nativo"

            return {
                "texto": _sanitizar(texto_final),
                "slides": total_slides,
                "metodo_extracao": metodo,
                "flags": [],
            }

    except Exception:
        return _falha("arquivo_ilegivel", slides=0)


def _extrair_text_frames(slide: Any) -> str:
    """Extrai texto de todos os text frames de um slide."""
    partes: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                linha = "".join(run.text for run in para.runs).strip()
                if linha:
                    partes.append(linha)
    return "\n".join(partes)


def _tem_imagem(shapes: Any) -> bool:
    """Verifica recursivamente se há alguma imagem nas shapes do slide. [R3]

    Canva exporta slides onde a imagem fica dentro de um GroupShape —
    verificação recursiva é necessária para detectar esses casos.
    """
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
        # R3 — GroupShape (Canva, Google Slides exportados)
        if hasattr(shape, "shapes") and _tem_imagem(shape.shapes):
            return True
    return False


def _ocr_slide(slide: Any) -> str:
    """Tenta extrair texto do slide via OCR heurístico.

    Como LibreOffice não está disponível no Streamlit Cloud, usa heurística:
    cria imagem em branco e extrai texto das shapes Picture encontradas.
    """
    textos_ocr: list[str] = []
    _ocr_shapes(slide.shapes, textos_ocr)
    return "\n".join(textos_ocr)


def _ocr_shapes(shapes: Any, acumulador: list[str]) -> None:
    """Percorre shapes recursivamente e aplica OCR nas imagens. [M3, R3]"""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_data = shape.image.blob
            img = Image.open(io.BytesIO(img_data))
            try:
                texto = pytesseract.image_to_string(img, lang="por")
                if texto.strip():
                    acumulador.append(texto.strip())
            finally:
                # M3 — garantir liberação de memória da imagem Pillow
                img.close()
        # R3 — descer em GroupShapes (Canva)
        elif hasattr(shape, "shapes"):
            _ocr_shapes(shape.shapes, acumulador)


def _sanitizar(texto: str) -> str:
    """Remove caracteres de controle e normaliza encoding UTF-8."""
    texto = unicodedata.normalize("NFC", texto)
    texto = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", texto)
    return texto.strip()


def _falha(flag: str, slides: int) -> dict[str, object]:
    return {
        "texto": "",
        "slides": slides,
        "metodo_extracao": "nenhum",
        "flags": [flag],
    }


def convert_pptx_batch(
    file_paths: list[Path], max_workers: int = MAX_WORKERS_BATCH
) -> list[dict[str, object]]:
    """Converte lista de arquivos PPTX em paralelo. [R5]

    Usa ThreadPoolExecutor — adequado para operações I/O-bound.
    Para batch com OCR pesado, reduzir max_workers para 2.
    """
    resultados: list[dict[str, object]] = [{}] * len(file_paths)
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = {executor.submit(convert_pptx, p): i for i, p in enumerate(file_paths)}
        for future in as_completed(futures):
            idx = futures[future]
            resultados[idx] = future.result()
    return resultados
