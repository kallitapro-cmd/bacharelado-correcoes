#!/usr/bin/env python3
"""
generate_fixtures.py — Geração determinística dos fixtures sintéticos para
o Corretor Acadêmico (Bacharelado em Administração).

Gera, no diretório `docs/fixtures/`:

1. `aba1-exemplo-anon.xlsx`  — planilha do PA com 3 abas (Alunos, Consolidado, Pós-Work Etapa 4)
2. `batch-exemplo/aluno-a-entrega.pdf`        — PDF acadêmico sintético
3. `batch-exemplo/aluno-b-apresentacao.pptx`  — PPTX com 3 slides
4. `batch-exemplo/aluno-c-relatorio.docx`     — DOCX com texto acadêmico

Todos os dados são **sintéticos e anonimizados**. O padrão de RA segue o ADR-001:
- 11 dígitos: padrão `2026010XXXX` ou `2025010XXXX`
- 10 dígitos: padrão `2026100XXX` (zero faltante na posição 4)
- caso especial: `0000100041` (fallback para Camada 3)

Uso:
    python3 docs/fixtures/generate_fixtures.py

Dependências (pip install --user):
    openpyxl, reportlab, python-pptx, python-docx
"""

from __future__ import annotations

from pathlib import Path

# Diretório base dos fixtures (mesmo diretório deste script)
FIXTURES_DIR = Path(__file__).resolve().parent
BATCH_DIR = FIXTURES_DIR / "batch-exemplo"

# ----------------------------------------------------------------------------
# Dados sintéticos compartilhados
# ----------------------------------------------------------------------------

# Mix representativo do padrão real da turma (124 alunos: 111 com 11 dígitos,
# 13 com 10 dígitos, 1 caso especial). Aqui levamos 8 amostras estratégicas.
ALUNOS: list[dict[str, str | float | int]] = [
    # 5 RAs de 11 dígitos (padrão canônico mais comum) -------------------------
    {
        "ra": "20260100418",
        "nome": "Aluno A",
        "email": "aluno.a@exemplo.edu.br",
        "telefone": "(11) 99999-0001",
        "turma": "ADM-2026",
        "grupo": "G1",
        "nota_a1": 8.5,
        "nota_a2": 9.0,
        "comentario": "Boa estruturação do raciocínio.",
    },
    {
        "ra": "20260100419",
        "nome": "Aluno B",
        "email": "aluno.b@exemplo.edu.br",
        "telefone": "(11) 99999-0002",
        "turma": "ADM-2026",
        "grupo": "G1",
        "nota_a1": 7.0,
        "nota_a2": 8.0,
        "comentario": "Precisa aprofundar a análise crítica.",
    },
    {
        "ra": "20260100501",
        "nome": "Aluno C",
        "email": "aluno.c@exemplo.edu.br",
        "telefone": "(11) 99999-0003",
        "turma": "ADM-2026",
        "grupo": "G2",
        "nota_a1": 9.5,
        "nota_a2": 9.0,
        "comentario": "Trabalho exemplar.",
    },
    {
        "ra": "20250100621",
        "nome": "Aluno D",
        "email": "aluno.d@exemplo.edu.br",
        "telefone": "(11) 99999-0004",
        "turma": "ADM-2025",
        "grupo": "G2",
        "nota_a1": 6.5,
        "nota_a2": 7.5,
        "comentario": "Bom desenvolvimento, revisar conclusão.",
    },
    {
        "ra": "20260100777",
        "nome": "Aluno E",
        "email": "aluno.e@exemplo.edu.br",
        "telefone": "(11) 99999-0005",
        "turma": "ADM-2026",
        "grupo": "G3",
        "nota_a1": 8.0,
        "nota_a2": 8.5,
        "comentario": "Boa apresentação.",
    },
    # 2 RAs de 10 dígitos (padrão legado, sem o zero da posição 4) -----------
    {
        "ra": "2026100420",
        "nome": "Aluno F",
        "email": "aluno.f@exemplo.edu.br",
        "telefone": "(11) 99999-0006",
        "turma": "ADM-2026",
        "grupo": "G3",
        "nota_a1": 7.5,
        "nota_a2": 8.0,
        "comentario": "Argumentação consistente.",
    },
    {
        "ra": "2025100333",
        "nome": "Aluno G",
        "email": "aluno.g@exemplo.edu.br",
        "telefone": "(11) 99999-0007",
        "turma": "ADM-2025",
        "grupo": "G4",
        "nota_a1": 6.0,
        "nota_a2": 7.0,
        "comentario": "Cumpriu o solicitado, sem destaques.",
    },
    # 1 RA caso especial (fallback Camada 3) ---------------------------------
    {
        "ra": "0000100041",
        "nome": "Aluno H",
        "email": "aluno.h@exemplo.edu.br",
        "telefone": "(11) 99999-0008",
        "turma": "ADM-2024",
        "grupo": "G4",
        "nota_a1": 8.0,
        "nota_a2": 7.5,
        "comentario": "Matrícula legada — revisar manualmente.",
    },
]

DATA_ETAPA = "27/04/2026"
ETAPA_LABEL = "Etapa 4"

# ----------------------------------------------------------------------------
# 1. Planilha XLSX (3 abas)
# ----------------------------------------------------------------------------


def generate_xlsx() -> Path:
    """Gera `aba1-exemplo-anon.xlsx` com 3 abas representando a planilha real do PA."""
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    wb = Workbook()

    # ---------- Aba 1: Alunos (Manifesto) ----------
    ws_alunos = wb.active
    ws_alunos.title = "Alunos"

    headers_alunos = ["RA", "Nome Completo", "E-mail", "Telefone"]
    bold = Font(bold=True)
    header_fill = PatternFill(start_color="DDDDDD", end_color="DDDDDD", fill_type="solid")

    for col_idx, header in enumerate(headers_alunos, start=1):
        cell = ws_alunos.cell(row=1, column=col_idx, value=header)
        cell.font = bold
        cell.fill = header_fill

    for row_idx, aluno in enumerate(ALUNOS, start=2):
        ws_alunos.cell(row=row_idx, column=1, value=aluno["ra"])
        ws_alunos.cell(row=row_idx, column=2, value=aluno["nome"])
        ws_alunos.cell(row=row_idx, column=3, value=aluno["email"])
        ws_alunos.cell(row=row_idx, column=4, value=aluno["telefone"])

    # Larguras de coluna razoáveis
    for col_idx, width in enumerate([14, 22, 28, 18], start=1):
        ws_alunos.column_dimensions[get_column_letter(col_idx)].width = width

    # ---------- Aba 2: Consolidado (notas por etapa) ----------
    ws_cons = wb.create_sheet(title="Consolidado")

    # Linha 1: cabeçalhos fixos (RA, Nome) + data mesclada cobrindo C1:E1
    # (A1 e B1 são células únicas — NÃO são mescladas com a data)
    ws_cons.cell(row=1, column=1, value="RA").font = bold
    ws_cons.cell(row=1, column=2, value="Nome").font = bold
    ws_cons.cell(row=1, column=3, value=DATA_ETAPA).font = bold
    ws_cons.merge_cells(start_row=1, start_column=3, end_row=1, end_column=5)
    ws_cons.cell(row=1, column=3).alignment = Alignment(horizontal="center")

    # Linha 2: A/B vazios (continuação dos cabeçalhos fixos), C-E com A1, A2, A3
    ws_cons.cell(row=2, column=1, value=None).fill = header_fill
    ws_cons.cell(row=2, column=2, value=None).fill = header_fill
    for idx, atividade in enumerate(["A1", "A2", "A3"], start=3):
        cell = ws_cons.cell(row=2, column=idx, value=atividade)
        cell.font = bold
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center")

    # Linhas 3+: RA + Nome + 3 notas sintéticas (A3 derivado da média)
    for row_idx, aluno in enumerate(ALUNOS, start=3):
        ws_cons.cell(row=row_idx, column=1, value=aluno["ra"])
        ws_cons.cell(row=row_idx, column=2, value=aluno["nome"])
        ws_cons.cell(row=row_idx, column=3, value=aluno["nota_a1"])
        ws_cons.cell(row=row_idx, column=4, value=aluno["nota_a2"])
        nota_a3 = round((float(aluno["nota_a1"]) + float(aluno["nota_a2"])) / 2, 1)
        ws_cons.cell(row=row_idx, column=5, value=nota_a3)

    for col_idx, width in enumerate([14, 22, 8, 8, 8], start=1):
        ws_cons.column_dimensions[get_column_letter(col_idx)].width = width

    # ---------- Aba 3: Pós-Work Etapa 4 (detalhe da etapa) ----------
    ws_etapa = wb.create_sheet(title="Pós-Work Etapa 4")

    # Linha 1: data da etapa (mesclada cobrindo as 6 colunas)
    ws_etapa.cell(row=1, column=1, value=DATA_ETAPA).font = bold
    ws_etapa.merge_cells(start_row=1, start_column=1, end_row=1, end_column=6)
    ws_etapa.cell(row=1, column=1).alignment = Alignment(horizontal="center")

    # Linha 2: rótulos das colunas
    headers_etapa = ["RA", "Nome", "Nota A1", "Nota A2", "Comentário", "Grupo"]
    for col_idx, header in enumerate(headers_etapa, start=1):
        cell = ws_etapa.cell(row=2, column=col_idx, value=header)
        cell.font = bold
        cell.fill = header_fill

    # Linhas 3+: dados dos alunos
    for row_idx, aluno in enumerate(ALUNOS, start=3):
        ws_etapa.cell(row=row_idx, column=1, value=aluno["ra"])
        ws_etapa.cell(row=row_idx, column=2, value=aluno["nome"])
        ws_etapa.cell(row=row_idx, column=3, value=aluno["nota_a1"])
        ws_etapa.cell(row=row_idx, column=4, value=aluno["nota_a2"])
        ws_etapa.cell(row=row_idx, column=5, value=aluno["comentario"])
        ws_etapa.cell(row=row_idx, column=6, value=aluno["grupo"])

    for col_idx, width in enumerate([14, 18, 10, 10, 40, 10], start=1):
        ws_etapa.column_dimensions[get_column_letter(col_idx)].width = width

    output_path = FIXTURES_DIR / "aba1-exemplo-anon.xlsx"
    wb.save(output_path)
    return output_path


# ----------------------------------------------------------------------------
# 2. PDF (entrega do Aluno A)
# ----------------------------------------------------------------------------


def generate_pdf() -> Path:
    """Gera um PDF acadêmico sintético para `batch-exemplo/aluno-a-entrega.pdf`."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    output_path = BATCH_DIR / "aluno-a-entrega.pdf"
    doc = SimpleDocTemplate(
        str(output_path),
        pagesize=A4,
        leftMargin=2 * cm,
        rightMargin=2 * cm,
        topMargin=2 * cm,
        bottomMargin=2 * cm,
        title="Aluno A — Entrega Etapa 4",
        author="Aluno A (sintético)",
    )
    styles = getSampleStyleSheet()
    story = []

    story.append(Paragraph("<b>RA:</b> 20260100418", styles["Normal"]))
    story.append(Paragraph("<b>Aluno:</b> Aluno A", styles["Normal"]))
    story.append(Paragraph("<b>Disciplina:</b> Projeto Extensionista de IA", styles["Normal"]))
    story.append(Paragraph(f"<b>Atividade:</b> {ETAPA_LABEL} — A1", styles["Normal"]))
    story.append(Spacer(1, 0.5 * cm))

    story.append(Paragraph("Introdução", styles["Heading2"]))
    story.append(
        Paragraph(
            "Este trabalho apresenta uma análise sintética do projeto desenvolvido durante a "
            "Etapa 4 do componente curricular. O objetivo é demonstrar a aplicação dos conceitos "
            "discutidos em sala, com foco em metodologias administrativas e governança de dados.",
            styles["BodyText"],
        )
    )
    story.append(Spacer(1, 0.3 * cm))

    story.append(Paragraph("Desenvolvimento", styles["Heading2"]))
    story.append(
        Paragraph(
            "Durante o desenvolvimento, foram aplicados frameworks de planejamento estratégico, "
            "considerando o contexto organizacional e os principais stakeholders. A análise SWOT "
            "preliminar identificou pontos fortes na cultura colaborativa e oportunidades de "
            "expansão no segmento de pequenas e médias empresas locais.",
            styles["BodyText"],
        )
    )
    story.append(Spacer(1, 0.3 * cm))

    story.append(
        Paragraph(
            "A metodologia adotada combinou pesquisa documental e entrevistas estruturadas com "
            "três especialistas. Os dados coletados foram organizados em planilhas temáticas, "
            "permitindo o cruzamento de informações qualitativas e quantitativas.",
            styles["BodyText"],
        )
    )
    story.append(Spacer(1, 0.3 * cm))

    story.append(Paragraph("Resultados", styles["Heading2"]))
    story.append(
        Paragraph(
            "Os resultados indicam que a integração entre as áreas operacional e estratégica "
            "produz ganhos mensuráveis em produtividade. Foram observados aumentos de eficiência "
            "em três indicadores-chave, com redução de retrabalho e maior previsibilidade de prazos.",
            styles["BodyText"],
        )
    )
    story.append(Spacer(1, 0.3 * cm))

    story.append(Paragraph("Conclusão", styles["Heading2"]))
    story.append(
        Paragraph(
            "Conclui-se que a aplicação prática dos conceitos estudados, aliada a um processo "
            "estruturado de coleta e análise de dados, gera evidências consistentes para a tomada "
            "de decisão. Como próximos passos, recomenda-se a validação dos achados em uma amostra "
            "ampliada e a institucionalização das práticas testadas.",
            styles["BodyText"],
        )
    )

    doc.build(story)
    return output_path


# ----------------------------------------------------------------------------
# 3. PPTX (apresentação do Aluno B)
# ----------------------------------------------------------------------------


def generate_pptx() -> Path:
    """Gera um PPTX com 3 slides para `batch-exemplo/aluno-b-apresentacao.pptx`."""
    from pptx import Presentation
    from pptx.util import Pt

    output_path = BATCH_DIR / "aluno-b-apresentacao.pptx"
    prs = Presentation()

    # Slide 1: capa -----------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # title slide
    slide1.shapes.title.text = "Etapa 4 — Apresentação"
    if slide1.placeholders[1] is not None:
        subtitle = slide1.placeholders[1]
        subtitle.text = (
            "Aluno: Aluno B  |  RA: 20260100419\n"
            "Projeto Extensionista de IA — ADM-2026"
        )
        for para in subtitle.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(18)

    # Slide 2: objetivos ------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
    slide2.shapes.title.text = "Objetivos do Trabalho"
    body = slide2.placeholders[1].text_frame
    body.text = "Apresentar a análise estratégica desenvolvida durante a etapa"
    for item in [
        "Aplicar frameworks administrativos clássicos (SWOT, 5 forças)",
        "Coletar evidências em entrevistas semiestruturadas",
        "Propor recomendações práticas baseadas em dados",
    ]:
        p = body.add_paragraph()
        p.text = item
        p.level = 1

    # Slide 3: conclusões -----------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Conclusões e Próximos Passos"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "A integração operacional-estratégica gera ganhos mensuráveis"
    for item in [
        "Redução de retrabalho em 18% nos processos analisados",
        "Maior previsibilidade de prazos nas entregas",
        "Próximo passo: replicar metodologia em outras áreas",
    ]:
        p = body3.add_paragraph()
        p.text = item
        p.level = 1

    prs.save(output_path)
    return output_path


# ----------------------------------------------------------------------------
# 4. DOCX (relatório do Aluno C)
# ----------------------------------------------------------------------------


def generate_docx() -> Path:
    """Gera um DOCX acadêmico sintético para `batch-exemplo/aluno-c-relatorio.docx`."""
    from docx import Document
    from docx.shared import Pt

    output_path = BATCH_DIR / "aluno-c-relatorio.docx"
    doc = Document()

    # Estilo base ----------------------------------------------------------
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    # Cabeçalho identificador ---------------------------------------------
    doc.add_heading("Relatório — Etapa 4", level=1)
    doc.add_paragraph("RA: 20260100501")
    doc.add_paragraph("Aluno: Aluno C")
    doc.add_paragraph("Disciplina: Projeto Extensionista de IA")
    doc.add_paragraph(f"Data: {DATA_ETAPA}")

    # Conteúdo acadêmico ---------------------------------------------------
    doc.add_heading("1. Introdução", level=2)
    doc.add_paragraph(
        "Este relatório sintetiza o trabalho desenvolvido durante a Etapa 4 do componente "
        "curricular, com foco em metodologias de gestão e análise de processos. O escopo abrange "
        "três frentes principais: levantamento de requisitos, modelagem de processos e proposta "
        "de melhoria contínua."
    )

    doc.add_heading("2. Metodologia", level=2)
    doc.add_paragraph(
        "A pesquisa adotou abordagem qualitativa, combinando análise documental e entrevistas "
        "com cinco profissionais da área. Foi utilizada a técnica de análise temática para "
        "categorizar os achados e identificar padrões recorrentes."
    )

    doc.add_heading("3. Resultados", level=2)
    doc.add_paragraph(
        "Os resultados indicam três eixos de oportunidade: (i) automação de tarefas repetitivas, "
        "(ii) padronização da comunicação entre áreas, e (iii) revisão dos indicadores de "
        "desempenho atualmente monitorados. Cada um desses eixos demanda intervenções específicas "
        "e foi detalhado em seções subsequentes."
    )

    doc.add_heading("4. Considerações Finais", level=2)
    doc.add_paragraph(
        "O trabalho permitiu consolidar uma visão integrada do processo estudado, com evidências "
        "robustas para fundamentar recomendações. Como continuidade, propõe-se a validação dos "
        "achados em uma segunda rodada de entrevistas e a elaboração de um plano de ação detalhado."
    )

    doc.save(output_path)
    return output_path


# ----------------------------------------------------------------------------
# Orquestração
# ----------------------------------------------------------------------------


def main() -> None:
    """Gera todos os fixtures, criando diretórios conforme necessário."""
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    BATCH_DIR.mkdir(parents=True, exist_ok=True)

    print("Gerando fixtures sintéticos em:", FIXTURES_DIR)

    xlsx_path = generate_xlsx()
    print(f"  [OK] XLSX  -> {xlsx_path.relative_to(FIXTURES_DIR.parent.parent)}")

    pdf_path = generate_pdf()
    print(f"  [OK] PDF   -> {pdf_path.relative_to(FIXTURES_DIR.parent.parent)}")

    pptx_path = generate_pptx()
    print(f"  [OK] PPTX  -> {pptx_path.relative_to(FIXTURES_DIR.parent.parent)}")

    docx_path = generate_docx()
    print(f"  [OK] DOCX  -> {docx_path.relative_to(FIXTURES_DIR.parent.parent)}")

    print("\nFixtures gerados com sucesso.")


if __name__ == "__main__":
    main()
